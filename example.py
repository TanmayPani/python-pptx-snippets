import marimo

__generated_with = "0.20.2"
app = marimo.App(width="medium")

with app.setup:
    import os

    import marimo as mo

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    from snippets import add_movie, move_slide

    prs = Presentation()
    hidden_slides = []
    file_name = mo.notebook_location() / "output.pptx"
    if os.path.exists(file_name):
        os.remove(file_name)


# @app.cell
# def _():
#    from snippets import get_thumbnail_from_video
#
#    _video_path = mo.notebook_location() / "public" / "example_video.mp4"
#    print(get_thumbnail_from_video(_video_path))
#    return


@app.cell(hide_code=True)
def _():
    mo.md("""
    Created a `Presentation` instance in the setup cell (global state),
    will add one slide per cell below. Any hidden slides (e.g. slides
    containing fullscreen videos) will be appended to the hidden_slides
    list. They will be transferred to the end of the slide-deck before
    saving it in the last cell.

    Recommended to run this notebook as a script using `uv run main.py` or `python main.py` in an environment with the required dependencies installed (look at the `pyproject.toml` in this repository).
    """)
    return


@app.cell
def _():
    _slide = prs.slides.add_slide(prs.slide_layouts[0])
    _slide.shapes.title.text = "Hello, World!"
    _slide.placeholders[1].text = "python-pptx was here!"
    return


@app.cell
def _():
    _slide = prs.slides.add_slide(prs.slide_layouts[1])
    _left = _top = _width = _height = Inches(1.0)
    _shape = _slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _left, _top, _width, _height
    )
    return


@app.cell(hide_code=True)
def _():
    mo.md("""
    Here we add a movie to the next slide using our own `add_movie` defined in `snippets.py` which wraps the `add_movie` method of the `pptx.slide.Slide` instance returned by `prs.slides.add_slide` method

    If the `add_fullscreen` keyword argument is passed as `True` (default), then the `add_movie` method will add another slide which has the same video taking up the whole space and also add the clickable shapes to go back and forth between the current slide containing the original video and the fullscreen slide containing the fullscreen video, immitating the experience of having an icon toggling fullscreen mode on/off.
    """)
    return


@app.cell
def _():
    _slide = prs.slides.add_slide(prs.slide_layouts[6])
    _movie, _fs_movie_slide, _thumbnail = add_movie(
        prs,
        _slide,
        str(mo.notebook_location() / "public" / "example_video.mp4"),
        0.0,
        0.0,
        prs.slide_width / 2,
        prs.slide_height / 2,
        add_fullscreen=True,
    )
    hidden_slides.append(_fs_movie_slide)
    return


@app.cell(hide_code=True)
def _():
    mo.md(r"""
    #### Limitations :
    - During a slide show session (triggered by clicking on "Play"/"Play from current slide" option in your slides app), the fullscreen version will autoplay only once for a video, next time you go to the "fullscreen mode" of the same video, you will have to click on the play button in the video playback controls
    - If you use keynote (MS PowerPoint equivalent in MacOS), then it doesn't support clickable actions pointing to a hidden slide. The functions provided here work best with MS PowerPoint, but if you are restricted to Keynote, then set the keyword argument `hide_fullscreen_slide` to `False` (default: `True`). So the above code would look like:
       ```{python}
       _movie, _fs_movie_slide  = add_movie(
           prs,
           _slide,
           "./example_video.mp4",
           0.,
           0.,
           prs.slide_width/2,
           prs.slide_height/2,
           hide_fullscreen_slide=False,
       )
       ```
       So the fullscreen slides will be completely visible in the slide deck and you would be able to navigate to them like a normal slide. In that case it is highly recommended to use the snippet given in the [last cell](#last_cell) to move these slides to the end of the slide decl (just before saving the pptx)
    - Not tested for any open source presentation software (e.g. LibreOffice, OpenOffice etc.). LibreOffice doesn't show playback control bar in the video slide so it would be a lot messier to rewrite the logic to never use the play button on the user end.
    """)
    return


@app.cell
def _():
    _slide = prs.slides.add_slide(prs.slide_layouts[1])
    _slide.shapes.title.text = "Dummy slide"
    return


@app.cell(hide_code=True)
def _():
    mo.md(r"""
    <a id="last_cell"></a>
    The last cell will move all the extra slides (e.g. fullscreen video slides) to end of the slide deck, so that they don't interfere with the flow of the slides while viewing the `.pptx` output in PowerPoint/Keynote or in presentation mode (if using Keynote). We use the `move_slide` function defined in `snippets.py`
    """)
    return


@app.cell
def _():
    for _slide in hidden_slides:
        _from_index = prs.slides.index(_slide)
        move_slide(prs, _from_index, -1)

    prs.save(file_name)
    return


if __name__ == "__main__":
    app.run()
