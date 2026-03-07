import subprocess
from pathlib import Path

import pptx.util
import pptx.presentation
import pptx.slide
import pptx.shapes.picture
from lxml import etree

import pptx
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.action import PP_ACTION
from pptx.util import Inches
from pptx.oxml import parse_xml


def xpath(el: pptx.oxml.shapes.ShapeElement, query: str):
    """Utility to query an `pptx.shapes.Shape`'s xml tree."""
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    return etree.ElementBase.xpath(el, query, namespaces=nsmap)


def autoplay_media(media: pptx.shapes.picture.Movie) -> None:
    """
    Utility to autoplay a media (currently only video) upon on entering the slide containing it.

    Args:
        media: `Shape` object containing the video.
    """
    el_id = xpath(media.element, ".//p:cNvPr")[0].attrib["id"]
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), ".//p:cond")[0]
    cond.set("delay", "0")
    cond.set("evt", "onBegin")


def move_slide(
    pres: pptx.presentation.Presentation, from_index: int, to_index: int
) -> None:
    """Move slide at position `from_index` in presentation `pres` to `to_index`"""
    slides = list(pres.slides._sldIdLst)
    if to_index < 0:
        to_index = len(slides) + to_index
    pres.slides._sldIdLst.remove(slides[from_index])
    pres.slides._sldIdLst.insert(to_index, slides[from_index])


def get_thumbnail_from_video(movie_file: str, img_format: str = ".jpg") -> str:
    video_input_path = Path(movie_file).resolve()
    img_output_path = video_input_path.parent / (video_input_path.stem + img_format)
    subprocess.call(
        [
            "ffmpeg",
            "-i",
            video_input_path,
            "-ss",
            "00:00:00.000",
            "-vframes",
            "1",
            img_output_path,
            "-y",
        ],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return str(img_output_path)


def add_movie(
    pres: pptx.presentation.Presentation,
    slide: pptx.slide.Slide,
    movie_file: str,
    left: pptx.util.Length,
    top: pptx.util.Length,
    width: pptx.util.Length,
    height: pptx.util.Length,
    mime_type: str = "video/mp4",
    poster_frame_image: str | None = None,
    add_fullscreen: bool = True,
    hide_fullscreen_slide: bool = True,
) -> (
    pptx.shapes.picture.Movie
    | tuple[pptx.shapes.picture.Movie, pptx.slide.Slide, pptx.shapes.picture.Picture]
):
    """
    Wrapper around add_movie method of a `pptx.slide.Slide` instance to add movies with functionality to toggle fullscreen mode

    Args:
        pres: the presentaion instance which contains the slide instance
        slide: slide instance to which we add the movie
        movie_file: path to the movie .mp4 file
        left: X-coordinate of the movie frame's top-left corner
        top: Y-coordinate of the movie frame's top-left corner
        width: width of the movie frame
        height: height of the movie frame
        mime_type: input to the mime_type keyword argument of slide.add_movie method.
        poster_frame_image: input to the poster_frame_image keyword argument of slide.add_movie method.
        add_fullscreen: Whether to add fullscreen toggling feature.
        hide_fullscreen_slide: Whether to hide the extra fullscreen slide or not. Recommend setting True if using PowerPoint and False if using Keynote.

    Returns:
        If `add_fullscreen == True`, returns a tuple of `(movie_shape, fullscreen_movie_slide)` else just returns the `movie_shape`
    """

    if add_fullscreen:
        thn_img_path = (
            poster_frame_image
            if poster_frame_image is not None
            else get_thumbnail_from_video(movie_file)
        )
        thn_img = slide.shapes.add_picture(thn_img_path, left, top, width, height)

        fs_btn_w, fs_btn_h = Inches(1.5), Inches(0.5)
        fs_btn_left = left + width - fs_btn_w - Inches(0.2)
        fs_btn_top = top + height - fs_btn_h - Inches(0.2)

        fs_btn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, fs_btn_left, fs_btn_top, fs_btn_w, fs_btn_h
        )
        fs_btn.text = "Fullscreen"

        fs_movie_slide = pres.slides.add_slide(pres.slide_layouts[6])
        if hide_fullscreen_slide:
            fs_movie_slide.element.set("show", "0")

        movie = fs_movie_slide.shapes.add_movie(
            movie_file,
            0,
            0,
            pres.slide_width,
            pres.slide_height,
            mime_type="video/mp4",
            poster_frame_image=thn_img_path,
        )
        # movie.click_action.hyperlink.address = None
        autoplay_media(movie)
        fs_movie_slide.element.append(
            parse_xml(
                '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" advOnClick="0"/>'
            )
        )

        fs_btn.click_action.target_slide = fs_movie_slide

        fs_exit_btn = fs_movie_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            pres.slide_width - fs_btn_w - Inches(0.2),
            Inches(0.2),
            fs_btn_h,
            fs_btn_h,
        )
        fs_exit_btn.text = "X"
        fs_exit_btn.click_action.target_slide = slide

        return movie, fs_movie_slide, thn_img

    movie = slide.shapes.add_movie(
        movie_file,
        left,
        top,
        width,
        height,
        mime_type="video/mp4",
        poster_frame_image=None,
    )
    return movie
